from flask import Flask, render_template
from pymongo import MongoClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from matplotlib import pyplot as plt
import io

import matplotlib
matplotlib.use('Agg')

app = Flask(__name__)

# MongoDB connection
client = MongoClient('localhost', 27017)
db = client['srgw_data_store']
# collection = db['your_collection_name']

@app.route('/')
def generate_ppt():
    # Create a new PowerPoint presentation
    print("Testing")
    prs = Presentation()

    # Create the first slide with a bar chart
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    chart_data = [
        ('Category 1', 20),
        ('Category 2', 30),
        ('Category 3', 40),
    ]
    x_data, y_data = zip(*chart_data)
    plt.bar(x_data, y_data)
    plt.xlabel('Categories')
    plt.ylabel('Values')
    plt.title('Bar Chart')
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png')
    plt.close()

    plt.show()
    left = Inches(1)
    top = Inches(1)
    pic = slide1.shapes.add_picture(img_stream, left, top, height=Inches(4))
    img_stream.close()

    # # Create the second slide with a pie chart
    # slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    # chart_data = [
    #     ('Category A', 45),
    #     ('Category B', 30),
    #     ('Category C', 25),
    # ]
    # labels, sizes = zip(*chart_data)
    # plt.pie(sizes, labels=labels, autopct='%1.1f%%')
    # plt.title('Pie Chart')
    # img_stream = io.BytesIO()
    # plt.savefig(img_stream, format='png')
    # plt.close()
    # left = Inches(1)
    # top = Inches(1)
    # pic = slide2.shapes.add_picture(img_stream, left, top, height=Inches(4))
    # img_stream.close()

    # # Create the third slide with a grid of 3x3 cards
    # slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    # card_data = [
    #     ('Card 1', 'Value 1'),
    #     ('Card 2', 'Value 2'),
    #     ('Card 3', 'Value 3'),
    #     ('Card 4', 'Value 4'),
    #     ('Card 5', 'Value 5'),
    #     ('Card 6', 'Value 6'),
    #     ('Card 7', 'Value 7'),
    #     ('Card 8', 'Value 8'),
    #     ('Card 9', 'Value 9'),
    # ]
    # left = top = Inches(1)
    # card_width = Inches(2)
    # card_height = Inches(1)
    # for title, value in card_data:
    #     card = slide3.shapes.add_rectangle(left, top, card_width, card_height)
    #     left += card_width + Inches(0.1)
    #     if left + card_width > Inches(7):
    #         left = Inches(1)
    #         top += card_height + Inches(0.1)
    #     text = card.text = slide3.shapes.add_textbox(left, top, card_width, card_height).text_frame
    #     text.add_paragraph(title).font.bold = True
    #     text.add_paragraph(value)
    #     text.paragraphs[0].alignment = text.paragraphs[1].alignment = PP_ALIGN.CENTER

    # # Save the presentation to a byte stream
    # ppt_stream = io.BytesIO()
    # prs.save(ppt_stream)
    # ppt_stream.seek(0)

    return "ppt_stream.read()"

if __name__ == '__main__':
    app.run(port=4000,debug=True)